from pptx import Presentation
from pptx.util import Inches
import ipdb
import os

#########################################################################################
#
#
# Additional references
#
# https://python-pptx.readthedocs.org/en/latest/user/quickstart.html#add-picture-example
#
#
#########################################################################################

img_path = 'c:\\MyPNGFiles\\'

file_names = ['file1.png']


prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
left = top = Inches(0.1)
for file_name in file_names:
    print('adding file {0}'.format(file_name))
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(img_path + file_name,left,top,height=Inches(6),width=Inches(10))
prs.save('MyRepresntation.pptx')